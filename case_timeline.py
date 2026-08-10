"""
case_timeline.py

Reconstruye la linea de tiempo de un caso a partir de uno o varios
documentos (.pdf, .docx, .txt):

1. Recibe archivos sueltos y/o carpetas como argumentos.
2. Extrae el texto de cada documento y, por fragmentos, le pide al LLM
   configurado en config.py (mismo patron que contract_risk_analyzer.py)
   que identifique todos los eventos con fecha: que paso, quien participo,
   en que categoria cae, y como se relaciona con otros eventos cercanos.
3. Normaliza las fechas (dateparser, reutilizando la logica ya probada en
   deadline_extractor.py) y ordena los eventos cronologicamente.
4. Le pide al LLM (una sola llamada, no una por evento) un resumen
   narrativo del caso y un analisis de inconsistencias: fechas que no
   coinciden entre documentos, eventos sin fecha, orden ilogico, posibles
   eventos faltantes.
5. Genera hasta 5 archivos de salida: grafico PNG (matplotlib), HTML
   interactivo (Mermaid.js + filtros por categoria), reporte Markdown,
   presentacion PPTX (python-pptx) y un JSON con los datos crudos.

Uso:
    python case_timeline.py expediente.pdf contrato.docx correos/
    python case_timeline.py caso_contrato/
    python case_timeline.py --format html demanda.pdf contestacion.docx

Nota: las relaciones entre eventos se expresan como texto dentro de cada
evento (ej. "ocurrio 15 dias despues de la firma del contrato"); no se
genera un grafo de nodos aparte (decision confirmada -- los 4 formatos
pedidos son de linea de tiempo, no de grafo).
"""

import argparse
import html
import json
import re
import sys
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path

from docx import Document
from pypdf import PdfReader

import config
from deadline_extractor import fecha_le_falta_anio, normalizar_fecha, resolver_anio_ambiguo_cli
from legal_research import ask_llm

# La consola de Windows por defecto (cp1252) no puede imprimir los emojis
# usados en la salida de este script; se reconfigura a UTF-8 si el stream
# lo permite (mismo gotcha ya resuelto en deadline_extractor.py).
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass


# --- Categorias de eventos --------------------------------------------------

CATEGORIAS = {
    "actos_juridicos": "Actos jurídicos",
    "actos_procesales": "Actos procesales",
    "comunicaciones": "Comunicaciones",
    "pagos": "Pagos",
    "plazos": "Plazos",
    "audiencias": "Audiencias",
    "incumplimientos": "Incumplimientos",
}

CATEGORIA_EMOJI = {
    "actos_juridicos": "📜",
    "actos_procesales": "⚖️",
    "comunicaciones": "💬",
    "pagos": "💰",
    "plazos": "📅",
    "audiencias": "🏛️",
    "incumplimientos": "🚨",
}

# Paleta pedida por el cliente. Solo trajo 6 colores para 7 categorias (no
# incluyo uno para "Actos procesales", distinto de "Actos juridicos" que si
# tiene) -- se asume un azul grisaceo coherente con el resto de la paleta;
# ajustar si no convence.
CATEGORIA_COLOR = {
    "actos_juridicos": "#2E86AB",
    "actos_procesales": "#4C6A8C",  # asumido, no venia en la paleta original
    "comunicaciones": "#A23B72",
    "pagos": "#F18F01",
    "incumplimientos": "#C73E1D",
    "plazos": "#6A4C93",
    "audiencias": "#D95D39",
}

EXTENSIONES_SOPORTADAS = (".pdf", ".docx", ".txt")

# Igual que en contract_risk_analyzer.py: fragmentos grandes porque aqui no
# hay busqueda semantica, se recorre el documento completo en bloques.
CHUNK_SIZE_WORDS = 4000


# --- Estructura de datos -----------------------------------------------------

@dataclass
class Evento:
    fecha_cruda: str
    descripcion: str
    personas: list[str]
    categoria: str
    relacion: str
    fuente: str
    fecha: date | None = None
    anio_ambiguo: bool = False


# --- Entrada: archivos y carpetas -------------------------------------------

def expandir_rutas(paths: list[str]) -> list[Path]:
    """Convierte los argumentos (archivos y/o carpetas) en una lista de
    archivos concretos con extension soportada. Las carpetas se recorren
    solo en su nivel superior (no recursivo, por simplicidad)."""
    resultado = []
    for raw in paths:
        p = Path(raw)
        if p.is_dir():
            for hijo in sorted(p.iterdir()):
                if hijo.is_file() and hijo.suffix.lower() in EXTENSIONES_SOPORTADAS:
                    resultado.append(hijo)
        elif p.is_file():
            if p.suffix.lower() in EXTENSIONES_SOPORTADAS:
                resultado.append(p)
            else:
                print(f"Aviso: formato no soportado, se omite: {p}", file=sys.stderr)
        else:
            print(f"Aviso: no existe, se omite: {p}", file=sys.stderr)
    return resultado


def extract_text(path: Path) -> str:
    """Extrae el texto de un PDF, DOCX o TXT. Sin OCR (fuera de alcance,
    igual que deadline_extractor.py)."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        doc = Document(str(path))
        # Documentos legales suelen tener datos clave (fechas, montos, partes)
        # en tablas, no solo en parrafos -- doc.paragraphs por si solo las
        # ignora por completo, y un documento con contenido solo en tablas
        # terminaba extrayendo texto vacio.
        partes = [p.text for p in doc.paragraphs]
        for tabla in doc.tables:
            for fila in tabla.rows:
                for celda in fila.cells:
                    if celda.text.strip():
                        partes.append(celda.text)
        return "\n".join(partes)
    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")
    raise ValueError(f"Formato no soportado: '{suffix}'. Usa un archivo .pdf, .docx o .txt")


def chunk_text(text: str, chunk_size_words: int = CHUNK_SIZE_WORDS) -> list[str]:
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size_words):
        chunk = " ".join(words[i:i + chunk_size_words])
        if chunk.strip():
            chunks.append(chunk)
    return chunks or [""]


# --- Extraccion de eventos via LLM ------------------------------------------

def build_extraction_prompt(chunk: str) -> str:
    catalogo = "\n".join(f"- {CATEGORIAS[k]} ({k})" for k in CATEGORIAS)
    return (
        "Eres un abogado especializado en reconstruir la cronologia de un caso. "
        "Lee EXCLUSIVAMENTE el siguiente fragmento de un documento del expediente "
        "y extrae TODOS los eventos con fecha que menciona (hechos, actos "
        "juridicos o procesales, comunicaciones, pagos, plazos, audiencias, "
        "incumplimientos).\n\n"
        f"Categorias disponibles:\n{catalogo}\n\n"
        "Responde UNICAMENTE con un arreglo JSON (sin texto adicional, sin "
        "bloques de codigo markdown), donde cada elemento representa un evento "
        "y tiene esta forma exacta:\n"
        '[{"fecha": "fecha tal como aparece en el texto, o cadena vacia si el '
        'evento no tiene fecha explicita", "descripcion": "que paso, en una '
        'oracion breve", "personas": ["nombre 1", "nombre 2"], "categoria": '
        '"una de las claves del catalogo anterior, ej. actos_juridicos", '
        '"relacion": "como se relaciona con otro evento cercano, ej. \'ocurrio '
        "10 dias despues de la firma del contrato', o cadena vacia si no "
        'aplica"}]\n\n'
        "Si el fragmento no contiene ningun evento con fecha relevante, "
        "responde exactamente: []\n\n"
        f"Fragmento del documento:\n\n{chunk}"
    )


def parse_llm_events(raw: str) -> list[dict]:
    cleaned = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.IGNORECASE | re.MULTILINE).strip()
    match = re.search(r"\[.*\]", cleaned, flags=re.DOTALL)
    if not match:
        return []
    try:
        data = json.loads(match.group(0))
    except json.JSONDecodeError:
        return []
    if not isinstance(data, list):
        return []

    eventos = []
    for item in data:
        if not isinstance(item, dict):
            continue
        categoria = str(item.get("categoria", "")).strip().lower()
        if categoria not in CATEGORIAS:
            continue
        descripcion = str(item.get("descripcion", "")).strip()
        if not descripcion:
            continue
        personas_raw = item.get("personas", [])
        personas = (
            [str(p).strip() for p in personas_raw if str(p).strip()]
            if isinstance(personas_raw, list) else []
        )
        eventos.append({
            "fecha_cruda": str(item.get("fecha", "")).strip(),
            "descripcion": descripcion,
            "personas": personas,
            "categoria": categoria,
            "relacion": str(item.get("relacion", "")).strip(),
        })
    return eventos


def extraer_eventos_documento(texto: str, filename: str, provider_config: dict) -> list[Evento]:
    chunks = chunk_text(texto)
    eventos: list[Evento] = []
    for idx, chunk in enumerate(chunks, start=1):
        if len(chunks) > 1:
            print(f"  Analizando fragmento {idx}/{len(chunks)} de {filename}...")
        prompt = build_extraction_prompt(chunk)
        try:
            raw = ask_llm(prompt, provider_config)
        except Exception as exc:
            print(f"  Aviso: fallo el analisis del fragmento {idx} de {filename}: {exc}", file=sys.stderr)
            continue
        for item in parse_llm_events(raw):
            eventos.append(Evento(
                fecha_cruda=item["fecha_cruda"],
                descripcion=item["descripcion"],
                personas=item["personas"],
                categoria=item["categoria"],
                relacion=item["relacion"],
                fuente=filename,
            ))
    return eventos


# --- Normalizacion de fechas y orden -----------------------------------------

def resolver_fechas(eventos: list[Evento], interactivo: bool = True) -> None:
    """Normaliza fecha_cruda -> fecha (date) en cada Evento, in-place. Los
    eventos sin ningun texto de fecha quedan con fecha=None de forma
    permanente (se resuelven mas adelante, en analizar_consistencia, que le
    pide al LLM sugerir una fecha probable por contexto)."""
    for e in eventos:
        texto = e.fecha_cruda.strip()
        if not texto:
            continue
        e.anio_ambiguo = fecha_le_falta_anio(texto)
        anio_hint = None
        if e.anio_ambiguo:
            anio_hint = (
                resolver_anio_ambiguo_cli(texto, e.descripcion)
                if interactivo else date.today().year
            )
        e.fecha = normalizar_fecha(texto, anio_hint=anio_hint)


def ordenar_eventos(eventos: list[Evento]) -> list[Evento]:
    con_fecha = [e for e in eventos if e.fecha is not None]
    sin_fecha = [e for e in eventos if e.fecha is None]
    con_fecha.sort(key=lambda e: e.fecha)
    return con_fecha + sin_fecha


# --- Analisis de consistencia y narrativa (una sola llamada al LLM) ---------

def build_consistency_prompt(eventos_ordenados: list[Evento]) -> str:
    lineas = []
    for idx, e in enumerate(eventos_ordenados, start=1):
        fecha_txt = e.fecha.strftime("%Y-%m-%d") if e.fecha else "SIN FECHA"
        lineas.append(
            f"{idx}. [{fecha_txt}] ({CATEGORIAS[e.categoria]}) {e.descripcion} "
            f"-- fuente: {e.fuente}"
        )
    lista = "\n".join(lineas)
    return (
        "Eres un abogado que reconstruye la cronologia de un caso a partir de "
        "una lista de eventos ya extraidos de varios documentos del "
        "expediente. Analiza la siguiente lista (ya ordenada cronologicamente; "
        "los eventos sin fecha detectada van al final):\n\n"
        f"{lista}\n\n"
        "Responde UNICAMENTE con un objeto JSON (sin texto adicional, sin "
        "bloques de codigo markdown) con esta forma exacta:\n"
        '{"narrativa": "texto corrido en espanol, en orden cronologico, '
        "contando la historia del caso (ej. 'El 15 de enero de 2024, las "
        "partes firmaron...')\", "
        '"inconsistencias": [{"tipo": '
        '"fecha_conflicto|sin_fecha|orden_ilogico|posible_faltante", '
        '"descripcion": "explicacion breve", "fecha_sugerida": "YYYY-MM-DD o '
        'cadena vacia"}]}\n\n'
        "Senala especificamente: eventos que parecen el mismo hecho citado en "
        "mas de un documento con fechas distintas, eventos sin fecha "
        "(sugiriendo una fecha probable por contexto si es posible), eventos "
        "que parecen fuera de orden logico, y posibles eventos faltantes que "
        "se infieren del contexto pero no aparecen explicitamente en la lista."
    )


def analizar_consistencia(eventos_ordenados: list[Evento], provider_config: dict) -> dict:
    if not eventos_ordenados:
        return {"narrativa": "", "inconsistencias": []}
    prompt = build_consistency_prompt(eventos_ordenados)
    try:
        raw = ask_llm(prompt, provider_config)
    except Exception as exc:
        print(f"Aviso: fallo el analisis de consistencia: {exc}", file=sys.stderr)
        return {"narrativa": "", "inconsistencias": []}

    cleaned = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.IGNORECASE | re.MULTILINE).strip()
    match = re.search(r"\{.*\}", cleaned, flags=re.DOTALL)
    if not match:
        return {"narrativa": "", "inconsistencias": []}
    try:
        data = json.loads(match.group(0))
    except json.JSONDecodeError:
        return {"narrativa": "", "inconsistencias": []}

    narrativa = str(data.get("narrativa", "")).strip()
    inconsistencias = []
    inconsistencias_raw = data.get("inconsistencias", [])
    if isinstance(inconsistencias_raw, list):
        for item in inconsistencias_raw:
            if not isinstance(item, dict):
                continue
            inconsistencias.append({
                "tipo": str(item.get("tipo", "")).strip() or "otro",
                "descripcion": str(item.get("descripcion", "")).strip(),
                "fecha_sugerida": str(item.get("fecha_sugerida", "")).strip(),
            })
    return {"narrativa": narrativa, "inconsistencias": inconsistencias}


# --- Estructura de datos consolidada (una sola fuente de verdad) -----------

def construir_datos_timeline(nombre_caso: str, eventos_ordenados: list[Evento], analisis: dict) -> dict:
    con_fecha = [e for e in eventos_ordenados if e.fecha is not None]
    periodo_inicio = con_fecha[0].fecha.isoformat() if con_fecha else None
    periodo_fin = con_fecha[-1].fecha.isoformat() if con_fecha else None

    conteo = {k: 0 for k in CATEGORIAS}
    for e in eventos_ordenados:
        conteo[e.categoria] += 1

    return {
        "caso": nombre_caso,
        "generado_en": datetime.now().isoformat(timespec="seconds"),
        "periodo": {"inicio": periodo_inicio, "fin": periodo_fin},
        "eventos": [
            {
                "fecha": e.fecha.isoformat() if e.fecha else None,
                "fecha_original": e.fecha_cruda,
                "descripcion": e.descripcion,
                "personas": e.personas,
                "categoria": e.categoria,
                "categoria_label": CATEGORIAS[e.categoria],
                "fuente": e.fuente,
                "relacion": e.relacion,
            }
            for e in eventos_ordenados
        ],
        "estadisticas_por_categoria": {CATEGORIAS[k]: v for k, v in conteo.items()},
        "narrativa": analisis.get("narrativa", ""),
        "inconsistencias": analisis.get("inconsistencias", []),
    }


# --- Generacion de salidas ----------------------------------------------------

def generar_json(datos: dict) -> str:
    return json.dumps(datos, ensure_ascii=False, indent=2)


def generar_reporte_markdown(datos: dict) -> str:
    lines = []
    lines.append(f"# Línea de Tiempo del Caso: {datos['caso']}")
    lines.append("")
    lines.append(f"**Fecha de generación:** {datos['generado_en']}")
    lines.append("")

    eventos = datos["eventos"]
    n_docs = len({e["fuente"] for e in eventos}) if eventos else 0

    lines.append("## 📋 Resumen ejecutivo")
    lines.append("")
    lines.append(f"Se encontraron **{len(eventos)}** eventos en **{n_docs}** documento(s).")
    if datos["periodo"]["inicio"] and datos["periodo"]["fin"]:
        lines.append(f"Periodo analizado: **{datos['periodo']['inicio']}** a **{datos['periodo']['fin']}**.")
    lines.append("")

    lines.append("## 📊 Estadísticas por categoría")
    lines.append("")
    for label, count in datos["estadisticas_por_categoria"].items():
        if count:
            lines.append(f"- {label}: {count}")
    lines.append("")

    lines.append("## 🕒 Tabla de eventos")
    lines.append("")
    lines.append("| # | Fecha | Categoría | Descripción | Personas | Fuente |")
    lines.append("|---|---|---|---|---|---|")
    for idx, e in enumerate(eventos, start=1):
        fecha_txt = e["fecha"] or "_sin fecha_"
        personas_txt = ", ".join(e["personas"]) or "—"
        descripcion = e["descripcion"].replace("|", "/")
        lines.append(
            f"| {idx} | {fecha_txt} | {e['categoria_label']} | {descripcion} "
            f"| {personas_txt} | {e['fuente']} |"
        )
    lines.append("")

    if datos["narrativa"]:
        lines.append("## 📖 Narrativa del caso")
        lines.append("")
        lines.append(datos["narrativa"])
        lines.append("")

    lines.append("## 🔍 Inconsistencias encontradas")
    lines.append("")
    if datos["inconsistencias"]:
        for inc in datos["inconsistencias"]:
            sugerida = f" (fecha sugerida: {inc['fecha_sugerida']})" if inc.get("fecha_sugerida") else ""
            lines.append(f"- **{inc['tipo']}**: {inc['descripcion']}{sugerida}")
    else:
        lines.append("_No se detectaron inconsistencias relevantes._")
    lines.append("")

    return "\n".join(lines)


def generar_grafico_png(datos: dict) -> bytes:
    """Grafico de barras horizontal (una fila por evento, orden cronologico),
    coloreado por categoria. Los eventos sin fecha resuelta no se pueden
    ubicar en el eje X y se excluyen del grafico (se listan en el reporte
    y el JSON de todas formas)."""
    import io as _io

    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.dates as mdates
    import matplotlib.pyplot as plt

    plt.rcParams.update({"font.family": "sans-serif", "font.sans-serif": ["Arial", "DejaVu Sans"]})

    eventos_con_fecha = [e for e in datos["eventos"] if e["fecha"]]

    if not eventos_con_fecha:
        fig, ax = plt.subplots(figsize=(10, 2), facecolor="white")
        ax.text(0.5, 0.5, "No hay eventos con fecha para graficar", ha="center", va="center")
        ax.axis("off")
    else:
        fechas = [date.fromisoformat(e["fecha"]) for e in eventos_con_fecha]
        n = len(eventos_con_fecha)
        fig, ax = plt.subplots(figsize=(12, max(4, n * 0.5)), facecolor="white")

        # El ancho de la barra se escala con el rango total de fechas (no un
        # valor fijo de dias): con un rango de meses una barra de "2 dias" es
        # invisible, y con un rango de anios (como aqui, por el evento sin
        # anio que cayo en el actual) un ancho fijo mas grande se veria como
        # una barra de duracion en vez de un marcador puntual.
        rango_dias = max((max(fechas) - min(fechas)).days, 1)
        ancho_barra = max(rango_dias * 0.006, 1)

        for i, (e, fecha) in enumerate(zip(eventos_con_fecha, fechas)):
            color = CATEGORIA_COLOR.get(e["categoria"], "#666666")
            x = mdates.date2num(fecha)
            ax.barh(i, ancho_barra, left=x, height=0.5, color=color, edgecolor="none")

        # Fecha + descripcion como una sola etiqueta del eje Y (un solo
        # tamano de letra, en vez de dos columnas de texto dibujadas a mano):
        # matplotlib mide el ancho real de cada etiqueta y reserva el margen
        # izquierdo que haga falta, sin importar cuan larga sea. Dos columnas
        # de texto con anchos de columna fijos (intento anterior) se
        # superponian en cuanto una descripcion era mas larga que el ancho
        # de columna asumido -- este enfoque no tiene ese problema porque no
        # asume ningun ancho fijo.
        ax.set_yticks(range(n))
        ax.set_yticklabels(
            [f"{fecha.strftime('%d/%m/%Y')}   {e['descripcion'][:60]}"
             for e, fecha in zip(eventos_con_fecha, fechas)],
            fontsize=13,
        )
        ax.invert_yaxis()  # cronologico de arriba a abajo
        ax.xaxis_date()
        ax.xaxis.set_major_formatter(mdates.DateFormatter("%b %Y"))
        ax.set_xlabel("Fecha")
        ax.set_title(datos["caso"], fontsize=16, fontweight="bold")
        ax.grid(axis="x", linestyle="--", alpha=0.3)
        for spine in ("top", "right", "left"):
            ax.spines[spine].set_visible(False)

        categorias_presentes = sorted({e["categoria"] for e in eventos_con_fecha}, key=list(CATEGORIAS).index)
        handles = [plt.Rectangle((0, 0), 1, 1, color=CATEGORIA_COLOR[c]) for c in categorias_presentes]
        labels = [CATEGORIAS[c] for c in categorias_presentes]
        ax.legend(handles, labels, loc="upper left", bbox_to_anchor=(1.01, 1), fontsize=10, frameon=False)

    buf = _io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


def _mermaid_task_id(idx: int) -> str:
    return f"ev{idx}"


def generar_html_mermaid(datos: dict) -> str:
    """HTML autocontenido con un diagrama Mermaid (gantt) de la linea de
    tiempo, filtros por categoria (reconstruyen la definicion Mermaid en JS
    y la vuelven a renderizar, en vez de manipular el SVG ya generado) y
    tooltips nativos (<title> SVG) con los detalles de cada evento."""
    eventos = datos["eventos"]

    fecha_fallback = next((e["fecha"] for e in eventos if e["fecha"]), None) or date.today().isoformat()

    secciones: dict[str, list[tuple[int, dict]]] = {}
    for idx, e in enumerate(eventos, start=1):
        secciones.setdefault(e["categoria_label"], []).append((idx, e))

    lineas_mermaid = [
        "gantt",
        "    dateFormat  YYYY-MM-DD",
        "    axisFormat  %d/%m/%y",
        f"    title {datos['caso']}",
    ]
    for label, items in secciones.items():
        lineas_mermaid.append(f"    section {label}")
        for idx, e in items:
            tid = _mermaid_task_id(idx)
            fecha_iso = e["fecha"] or fecha_fallback
            # ":" y "," son estructuralmente significativos en la sintaxis de
            # tareas de mermaid gantt -- se sustituyen para no romper el parseo.
            titulo = e["descripcion"].replace(":", " -").replace(",", ";")[:70]
            lineas_mermaid.append(f"    {titulo} :{tid}, {fecha_iso}, 1d")
    mermaid_def = "\n".join(lineas_mermaid)

    tasks_meta = [
        {
            "id": _mermaid_task_id(idx),
            "categoria": e["categoria"],
            "fecha": e["fecha"] or "Sin fecha determinada",
            "descripcion": e["descripcion"],
            "personas": ", ".join(e["personas"]) or "—",
            "fuente": e["fuente"],
        }
        for idx, e in enumerate(eventos, start=1)
    ]
    categorias_presentes = list(dict.fromkeys(e["categoria"] for e in eventos))

    def _json_seguro(obj) -> str:
        # Evita que un "</script" dentro de un texto libre (descripcion del
        # LLM) cierre el bloque <script> antes de tiempo.
        return json.dumps(obj, ensure_ascii=False).replace("</", "<\\/")

    tasks_json = _json_seguro(tasks_meta)
    categorias_json = _json_seguro([{"key": k, "label": CATEGORIAS[k]} for k in categorias_presentes])
    mermaid_base_js = _json_seguro(mermaid_def)

    return f"""<!doctype html>
<html lang="es">
<head>
<meta charset="utf-8">
<title>Línea de tiempo: {html.escape(datos['caso'])}</title>
<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
<style>
  body {{ font-family: Arial, sans-serif; margin: 2rem; color: #211d1a; background: #fff; }}
  h1 {{ font-size: 1.4rem; }}
  #filtros {{ margin-bottom: 1rem; }}
  #filtros button {{
    margin: 0 0.4rem 0.4rem 0; padding: 0.35rem 0.9rem; border-radius: 999px;
    border: 1px solid #ccc; background: #f1ece1; cursor: pointer; font-size: 0.85rem;
  }}
  #filtros button.activo {{ background: #1d3a5f; color: #fff; border-color: #1d3a5f; }}
  #diagrama {{ overflow-x: auto; border: 1px solid #ddd; border-radius: 8px; padding: 1rem; }}
  .nota {{ color: #8c8477; font-size: 0.8rem; margin-top: 1rem; }}
</style>
</head>
<body>
<h1>Línea de tiempo: {html.escape(datos['caso'])}</h1>
<div id="filtros"></div>
<div id="diagrama"><pre class="mermaid" id="mermaid-src">{html.escape(mermaid_def)}</pre></div>
<p class="nota">Pasa el mouse sobre un evento para ver su detalle (personas, fuente, fecha).</p>

<script>
mermaid.initialize({{ startOnLoad: false }});

const TODOS_LOS_EVENTOS = {tasks_json};
const CATEGORIAS = {categorias_json};
const MERMAID_BASE = {mermaid_base_js};

let categoriaActiva = null; // null = todas

function construirDefinicion() {{
  if (!categoriaActiva) return MERMAID_BASE;
  const idsActivos = new Set(
    TODOS_LOS_EVENTOS.filter(e => e.categoria === categoriaActiva).map(e => e.id)
  );
  const lineas = MERMAID_BASE.split("\\n");
  const salida = [];
  for (const linea of lineas) {{
    const match = linea.match(/:(ev\\d+),/);
    if (match && !idsActivos.has(match[1])) continue;
    salida.push(linea);
  }}
  return salida.join("\\n");
}}

async function renderizar() {{
  const contenedor = document.getElementById("diagrama");
  const def = construirDefinicion();
  const {{ svg }} = await mermaid.render("grafico-timeline", def);
  contenedor.innerHTML = svg;

  TODOS_LOS_EVENTOS.forEach(ev => {{
    const rect = contenedor.querySelector('[id*="' + ev.id + '"] rect, rect[id*="' + ev.id + '"]');
    if (!rect) return;
    const title = document.createElementNS("http://www.w3.org/2000/svg", "title");
    title.textContent = ev.descripcion + "\\nFecha: " + ev.fecha + "\\nPersonas: " + ev.personas + "\\nFuente: " + ev.fuente;
    rect.appendChild(title);
  }});
}}

function construirFiltros() {{
  const cont = document.getElementById("filtros");
  const btnTodos = document.createElement("button");
  btnTodos.textContent = "Todas las categorías";
  btnTodos.className = "activo";
  btnTodos.onclick = () => seleccionar(null, btnTodos);
  cont.appendChild(btnTodos);

  CATEGORIAS.forEach(c => {{
    const btn = document.createElement("button");
    btn.textContent = c.label;
    btn.onclick = () => seleccionar(c.key, btn);
    cont.appendChild(btn);
  }});
}}

function seleccionar(categoria, btnClickeado) {{
  categoriaActiva = categoria;
  document.querySelectorAll("#filtros button").forEach(b => b.classList.remove("activo"));
  btnClickeado.classList.add("activo");
  renderizar();
}}

construirFiltros();
renderizar();
</script>
</body>
</html>"""


def generar_pptx(datos: dict, imagen_png: bytes) -> bytes:
    import io as _io

    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Diapositiva 1: portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = datos["caso"]
    slide.placeholders[1].text = f"Línea de tiempo del caso\nGenerado el {datos['generado_en']}"

    # Diapositiva 2: resumen ejecutivo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resumen ejecutivo"
    body = slide.placeholders[1].text_frame
    body.text = f"Total de eventos: {len(datos['eventos'])}"
    if datos["periodo"]["inicio"]:
        p = body.add_paragraph()
        p.text = f"Periodo: {datos['periodo']['inicio']} a {datos['periodo']['fin']}"
    for label, count in datos["estadisticas_por_categoria"].items():
        if count:
            p = body.add_paragraph()
            p.text = f"{label}: {count}"

    # Diapositiva 3: linea de tiempo visual (reutiliza el PNG ya generado)
    slide = prs.slides.add_slide(blank)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)).text_frame
    tf.text = "Línea de tiempo visual"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    slide.shapes.add_picture(_io.BytesIO(imagen_png), Inches(0.5), Inches(1.0), width=Inches(9))

    # Diapositiva(s) 4+: tabla de eventos, paginada de a 15 filas
    eventos = datos["eventos"]
    por_pagina = 15
    for inicio in range(0, max(len(eventos), 1), por_pagina):
        lote = eventos[inicio:inicio + por_pagina]
        slide = prs.slides.add_slide(blank)
        titulo = "Tabla de eventos" + (" (cont.)" if inicio > 0 else "")
        tf = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.5)).text_frame
        tf.text = titulo
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True

        filas = len(lote) + 1
        tabla = slide.shapes.add_table(
            filas, 4, Inches(0.4), Inches(0.9), Inches(9.2), Inches(0.4 * filas)
        ).table
        for c, encabezado in enumerate(["Fecha", "Categoría", "Descripción", "Fuente"]):
            tabla.cell(0, c).text = encabezado
        for r, e in enumerate(lote, start=1):
            tabla.cell(r, 0).text = e["fecha"] or "Sin fecha"
            tabla.cell(r, 1).text = e["categoria_label"]
            tabla.cell(r, 2).text = e["descripcion"][:120]
            tabla.cell(r, 3).text = e["fuente"]
        for r in range(filas):
            for c in range(4):
                tabla.cell(r, c).text_frame.paragraphs[0].font.size = Pt(11)

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# --- CLI ---------------------------------------------------------------------

def cargar_documentos(paths: list[Path]) -> tuple[dict[str, str], list[tuple[str, str]]]:
    """Devuelve (textos, fallos). fallos es una lista de (nombre_archivo,
    motivo) -- se propaga hasta la UI/CLI en vez de solo imprimirse a
    stderr, que en la app de Streamlit es invisible para quien la usa."""
    textos = {}
    fallos: list[tuple[str, str]] = []
    for path in paths:
        try:
            texto = extract_text(path)
        except ValueError as exc:
            fallos.append((path.name, str(exc)))
            continue
        except Exception as exc:
            fallos.append((path.name, f"error inesperado al leer el archivo: {exc}"))
            continue
        if not texto.strip():
            fallos.append((
                path.name,
                "no se pudo extraer texto (¿está vacío, dañado, o es un PDF "
                "escaneado sin capa de texto? no se hace OCR en esta herramienta)",
            ))
            continue
        textos[path.name] = texto
    return textos, fallos


def _nombre_caso(paths_originales: list[str]) -> str:
    """Si se paso una sola ruta (una carpeta o un solo archivo), se usa su
    nombre como nombre del caso; con varios archivos sueltos, no hay un
    nombre obvio y se usa uno generico."""
    if len(paths_originales) == 1:
        p = Path(paths_originales[0])
        return p.name if p.is_dir() else p.stem
    return "Caso judicial"


def _preguntar_formato() -> str:
    print("\n¿Qué formato quieres generar?")
    print("1. Todos (gráfico + HTML + reporte + PPTX)")
    print("2. Solo gráfico (PNG)")
    print("3. Solo HTML interactivo")
    print("4. Solo reporte (MD)")
    print("5. Solo PowerPoint")
    respuesta = input("Elige una opción (Enter para 1): ").strip()
    return {"1": "todos", "2": "png", "3": "html", "4": "md", "5": "pptx"}.get(respuesta, "todos")


def procesar_documentos(
    rutas: list[Path],
    provider_config: dict,
    interactivo: bool = True,
    on_documento_procesado=None,
) -> tuple[list[Evento], dict, list[tuple[str, str]]]:
    """Punto de entrada reutilizable (CLI y UI de Streamlit): extrae texto,
    eventos, resuelve fechas ambiguas, ordena y analiza consistencia.
    on_documento_procesado(nombre, n_eventos), si se pasa, se llama despues
    de procesar cada documento (para imprimir progreso o actualizar la UI).
    El tercer valor devuelto son los fallos de extraccion de texto (nombre,
    motivo) -- se devuelven en vez de solo imprimirse a stderr, para que la
    UI de Streamlit (donde stderr es invisible para quien la usa) tambien
    pueda mostrarlos."""
    textos, fallos = cargar_documentos(rutas)
    if not textos:
        detalle = "; ".join(f'{nombre}: {motivo}' for nombre, motivo in fallos)
        raise ValueError(
            f"No se pudo extraer texto de ningún documento. Detalle: {detalle}"
            if detalle else "No se pudo extraer texto de ningún documento."
        )

    todos_los_eventos: list[Evento] = []
    conteo_por_doc = {}
    for nombre, texto in textos.items():
        eventos_doc = extraer_eventos_documento(texto, nombre, provider_config)
        conteo_por_doc[nombre] = len(eventos_doc)
        todos_los_eventos.extend(eventos_doc)
        if on_documento_procesado:
            on_documento_procesado(nombre, len(eventos_doc))

    if not todos_los_eventos:
        raise ValueError("No se encontraron eventos con fecha en los documentos.")

    resolver_fechas(todos_los_eventos, interactivo=interactivo)
    eventos_ordenados = ordenar_eventos(todos_los_eventos)
    return eventos_ordenados, conteo_por_doc, fallos


def main():
    parser = argparse.ArgumentParser(
        description="Reconstruye la línea de tiempo de un caso a partir de varios documentos"
    )
    parser.add_argument("paths", nargs="+", help="Archivos (.pdf/.docx/.txt) y/o carpetas a analizar")
    parser.add_argument(
        "--format", choices=["todos", "png", "html", "md", "pptx"], default=None,
        help="Formato de salida (si se omite, se pregunta de forma interactiva)",
    )
    parser.add_argument(
        "--no-interactivo", action="store_true",
        help="No preguntar por terminal (años ambiguos o formato); usa 'todos' y el año actual",
    )
    parser.add_argument(
        "-o", "--output-prefix", default="timeline_caso",
        help="Prefijo de los archivos de salida (default: timeline_caso)",
    )
    args = parser.parse_args()

    interactivo = not args.no_interactivo

    rutas = expandir_rutas(args.paths)
    if not rutas:
        print("No se encontraron documentos válidos (.pdf/.docx/.txt) en las rutas dadas.", file=sys.stderr)
        sys.exit(1)

    provider_config = config.get_active_provider_config()
    print(f"Proveedor de LLM activo: {provider_config['provider']} ({provider_config['model']})")
    print(f"\n📂 Procesando {len(rutas)} documento(s)...")

    def _reportar_progreso(nombre, n_eventos):
        print(f"📄 {nombre}: {n_eventos} eventos encontrados")

    try:
        eventos_ordenados, _, fallos = procesar_documentos(
            rutas, provider_config, interactivo=interactivo,
            on_documento_procesado=_reportar_progreso,
        )
    except ValueError as exc:
        print(f"⚠️  {exc}", file=sys.stderr)
        sys.exit(1)

    for nombre, motivo in fallos:
        print(f"⚠️  {nombre}: {motivo}", file=sys.stderr)

    print(f"\n✅ Total: {len(eventos_ordenados)} eventos extraídos")

    print("\n=== ANÁLISIS DE EVENTOS ===")
    con_fecha = [e for e in eventos_ordenados if e.fecha]
    if con_fecha:
        print(f"Periodo: {con_fecha[0].fecha.strftime('%d/%m/%Y')} - {con_fecha[-1].fecha.strftime('%d/%m/%Y')}")
    print("Categorías:")
    conteo = {}
    for e in eventos_ordenados:
        conteo[e.categoria] = conteo.get(e.categoria, 0) + 1
    for k in CATEGORIAS:
        if conteo.get(k):
            print(f"  - {CATEGORIAS[k]}: {conteo[k]}")

    print("\n🔎 Analizando consistencia y generando narrativa...")
    analisis = analizar_consistencia(eventos_ordenados, provider_config)

    if analisis["inconsistencias"]:
        print("\n🔍 Inconsistencias encontradas:")
        for inc in analisis["inconsistencias"]:
            print(f"  - {inc['descripcion']}")

    nombre_caso = _nombre_caso(args.paths)
    datos = construir_datos_timeline(nombre_caso, eventos_ordenados, analisis)

    formato = args.format
    if formato is None:
        formato = _preguntar_formato() if (interactivo and sys.stdin.isatty()) else "todos"

    print()
    png_bytes = None
    if formato in ("todos", "png"):
        png_bytes = generar_grafico_png(datos)
        Path(f"{args.output_prefix}.png").write_bytes(png_bytes)
        print(f"📊 Generando gráfico... ✅ {args.output_prefix}.png")
    if formato in ("todos", "pptx"):
        if png_bytes is None:
            png_bytes = generar_grafico_png(datos)
        Path(f"{args.output_prefix}.pptx").write_bytes(generar_pptx(datos, png_bytes))
        print(f"📊 Generando PowerPoint... ✅ {args.output_prefix}.pptx")
    if formato in ("todos", "html"):
        Path(f"{args.output_prefix}.html").write_text(generar_html_mermaid(datos), encoding="utf-8")
        print(f"🌐 Generando HTML interactivo... ✅ {args.output_prefix}.html")
    if formato in ("todos", "md"):
        Path(f"{args.output_prefix}.md").write_text(generar_reporte_markdown(datos), encoding="utf-8")
        print(f"📝 Generando reporte... ✅ {args.output_prefix}.md")

    Path(f"{args.output_prefix}_events.json").write_text(generar_json(datos), encoding="utf-8")
    print(f"💾 Datos en JSON... ✅ {args.output_prefix}_events.json")

    print("\n✅ ¡Línea de tiempo completada!")


if __name__ == "__main__":
    main()
